"""Generate comparative study deck (run: pip install python-pptx && python build_team_deck.py)."""
from pptx import Presentation
from pptx.util import Inches, Pt

OUT = "Site_Plan_Completeness_Comparative_Study.pptx"


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullet_slide(prs, title, bullets, notes=None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(18)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def add_two_col_slide(prs, title, left_title, left_bullets, right_title, right_bullets):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    # Title
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.6))
    box.text_frame.text = title
    box.text_frame.paragraphs[0].font.size = Pt(28)
    box.text_frame.paragraphs[0].font.bold = True
    # Left column
    lx, ly, lw, lh = Inches(0.5), Inches(1.1), Inches(4.4), Inches(5.5)
    left = slide.shapes.add_textbox(lx, ly, lw, lh)
    tf = left.text_frame
    tf.text = left_title
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(20)
    for b in left_bullets:
        p = tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(14)
    # Right column
    rx = Inches(5.1)
    right = slide.shapes.add_textbox(rx, ly, lw, lh)
    tf2 = right.text_frame
    tf2.text = right_title
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.size = Pt(20)
    for b in right_bullets:
        p = tf2.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(14)


def add_placeholder_slide(prs, title, placeholder_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    slide.shapes.title.text = title
    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5))
    tf = box.text_frame
    tf.text = "\n".join(placeholder_lines)
    for p in tf.paragraphs:
        p.font.size = Pt(16)
        p.font.italic = True


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Site plan completeness check",
        "Comparative study: prompt-based GPT vs fine-tuned model\nTeam presentation — [Your name / date]",
    )

    add_bullet_slide(
        prs,
        "Objective",
        [
            "Detect on each plan sheet whether a North direction symbol and PE / AHJ-style stamps are present.",
            "Compare two approaches: (1) prompt-based inference with a general GPT vision model, (2) fine-tuned model trained on labeled plan pages.",
            "Same downstream use case; different cost, latency, and consistency trade-offs.",
        ],
    )

    add_two_col_slide(
        prs,
        "Approach 1 — Prompt-based (OpenAI GPT vision)",
        "What we send",
        [
            "Full-page image (rasterized PDF page).",
            "System + user prompts with explicit definitions:",
            "  • North Direction Symbol (what counts / what to exclude).",
            "  • PE stamp vs City/AHJ approval stamp rules.",
            "Few-shot examples: 1–N image+answer pairs in context so the model mimics labeling style.",
            "Instruction: output JSON e.g. {\"stamp\": bool, \"north_arrow\": bool}.",
        ],
        "Characteristics",
        [
            "No separate training job; change prompts without retraining.",
            "Depends on base model + prompt quality + few-shot coverage.",
            "Token cost scales with prompt length + images per request.",
            "Good for iteration and baselines before investing in fine-tuning.",
        ],
    )

    add_bullet_slide(
        prs,
        "Approach 2 — Fine-tuned model (current pipeline)",
        [
            "Base model: GPT-4o vision-capable base (e.g. gpt-4o-2024-08-06) via OpenAI supervised fine-tuning.",
            "Training format: JSONL chat examples — system prompt + user instruction + page image (base64 data URL) + assistant JSON labels.",
            "Technique: supervised fine-tuning (SFT) on image+text messages; optional validation file for held-out metrics.",
            "Inference: same message shape as training; model id = fine-tuned endpoint (ft:…).",
            "Goal: align model to our definitions and label distribution on real plan sets.",
        ],
    )

    add_bullet_slide(
        prs,
        "Training & validation data (high level)",
        [
            "Source: PDF plan sets → one PNG per page under images/<pdf_name>/.",
            "labels.json → PDF names + page_list + per-page north_arrow / PE Stamp flags → vision_train.jsonl.",
            "labels_val.json → disjoint PDFs/pages → vision_validation.jsonl (no duplicate rows vs train).",
            "Only PDFs listed in either label file are rasterized; one script builds both JSONL files.",
            "Fine-tune job: upload train file; optionally upload validation file for validation loss/metrics in dashboard.",
        ],
    )

    add_placeholder_slide(
        prs,
        "Results — Prompt-based (fill in)",
        [
            "[Metric / dataset name]",
            "Accuracy / precision / recall (stamp): __________",
            "Accuracy / precision / recall (north_arrow): __________",
            "Latency per page (avg): __________",
            "Cost per N pages: __________",
            "Qualitative notes (failure modes): __________",
        ],
    )

    add_placeholder_slide(
        prs,
        "Results — Fine-tuned model (fill in)",
        [
            "[Same metric / dataset for fair comparison]",
            "Accuracy / precision / recall (stamp): __________",
            "Accuracy / precision / recall (north_arrow): __________",
            "Latency per page (avg): __________",
            "Cost per N pages (inference): __________",
            "Qualitative notes: __________",
        ],
    )

    add_placeholder_slide(
        prs,
        "Side-by-side summary (fill in)",
        [
            "When prompt-based wins: __________",
            "When fine-tuned wins: __________",
            "Recommendation for production: __________",
        ],
    )

    add_bullet_slide(
        prs,
        "Thank you / Q&A",
        [
            "Appendix: repo scripts — 1_prepare_training_validation_data.py, 2_fine_tune.py, 3_predict.py.",
            "Questions?",
        ],
    )

    prs.save(OUT)
    print("Saved:", OUT)


if __name__ == "__main__":
    main()
